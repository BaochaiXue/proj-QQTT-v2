from __future__ import annotations

import argparse
import csv
import gc
import json
import math
import multiprocessing as mp
import queue
import shutil
import subprocess
import sys
import time
from dataclasses import asdict, dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

import cv2
import numpy as np


ROOT = Path(__file__).resolve().parents[2]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))


from data_process.depth_backends import (  # noqa: E402
    DEFAULT_FFS_REPO,
    FastFoundationStereoSingleEngineTensorRTRunner,
    FastFoundationStereoTensorRTRunner,
    resolve_tensorrt_engine_static_batch_size,
)
from data_process.depth_backends.geometry import align_depth_to_color  # noqa: E402
from data_process.visualization.calibration_io import load_calibration_transforms  # noqa: E402
from data_process.visualization.io_artifacts import write_image, write_json  # noqa: E402
from data_process.visualization.io_case import (  # noqa: E402
    depth_to_camera_points,
    get_frame_count,
    load_case_metadata,
    resolve_case_dir,
    transform_points,
)
from data_process.visualization.roi import crop_points_to_bounds  # noqa: E402
from data_process.visualization.triplet_video_compare import _render_open3d_hidden_window  # noqa: E402
from data_process.visualization.views import build_original_camera_view_configs  # noqa: E402
from data_process.visualization.workflows.masked_camera_view_compare import (  # noqa: E402
    _image_size_from_color_path,
    _mask_rgb_image,
    _scale_intrinsic_matrix,
)
from data_process.visualization.workflows.masked_pointcloud_compare import (  # noqa: E402
    MIN_MASKED_POINT_COUNT_FOR_FOCUS,
    PHYSTWIN_DATA_PROCESS_MASK_CONTRACT,
    _expand_bounds,
    filter_camera_clouds_with_pixel_masks,
    load_union_masks_for_camera_clouds,
    parse_text_prompts,
    refine_pixel_masks_with_phystwin_data_process_mask,
)


ROUND_CASE_REFS = (
    "static/ffs_30_static_round1_20260410_235202",
    "static/ffs_30_static_round2_20260414",
    "static/ffs_30_static_round3_20260414",
)
MODEL_NAMES = ("23-36-37", "20-26-39", "20-30-48")
SCALE_VALUES = (1.0, 0.75, 0.5)
VALID_ITERS_VALUES = (8, 4, 2)
ENGINE_NAMES = ("single_engine_fp32", "two_stage_fp16")
CAMERA_IDS = (0, 1, 2)
FRAME_IDX_VISUAL_DEFAULT = 10
MASK_PROMPT_DEFAULT = "stuffed animal"
DEFAULT_MAX_DISP = 192
DEFAULT_POINT_SIZE = 2.0
DEFAULT_LOOK_DISTANCE = 1.0
DEFAULT_ZOOM = 0.55
STATIC_MASK_FALLBACK_ROOTS = {
    "Round 1": ROOT / "data" / "static" / "masked_pointcloud_compare_round1_frame_0000_stuffed_animal" / "_generated_masks" / "ffs" / "sam31_masks",
    "Round 2": ROOT / "data" / "static" / "masked_pointcloud_compare_round2_frame_0000_stuffed_animal" / "_generated_masks" / "ffs" / "sam31_masks",
    "Round 3": ROOT / "data" / "static" / "masked_pointcloud_compare_round3_frame_0000_stuffed_animal" / "_generated_masks" / "ffs" / "sam31_masks",
}
ROUND_BENCHMARK_RESULT_TIMEOUT_S = 1800.0


@dataclass(frozen=True)
class ExperimentConfig:
    experiment_id: str
    engine: str
    model_name: str
    model_path: str
    scale: float
    valid_iters: int
    max_disp: int
    engine_height: int
    engine_width: int
    artifact_dir: str


@dataclass
class ExperimentFailure:
    experiment_id: str
    engine: str
    model_name: str
    scale: float
    valid_iters: int
    stage: str
    message: str
    log_path: str | None = None


@dataclass
class ReplayFrame:
    frame_idx: int
    left_image: np.ndarray
    right_image: np.ndarray


@dataclass(frozen=True)
class RoundCameraBenchmarkJob:
    config: ExperimentConfig
    ffs_repo: str
    round_label: str
    case_dir: str
    camera_idx: int
    frame_idx: int
    k_ir_left: np.ndarray
    k_color: np.ndarray
    t_ir_left_to_color: np.ndarray
    baseline_m: float
    color_output_shape: tuple[int, int]


@dataclass
class RoundCameraBenchmarkResult:
    round_label: str
    camera_idx: int
    fps: float
    depth_color_m: np.ndarray


@dataclass
class RoundCaseBundle:
    round_label: str
    case_ref: str
    case_dir: Path
    metadata: dict[str, Any]
    serial_numbers: list[str]
    k_ir_left_list: list[np.ndarray]
    k_color_list: list[np.ndarray]
    t_ir_left_to_color_list: list[np.ndarray]
    baselines_m: list[float]
    c2w_list: list[np.ndarray]
    stereo_pairs_by_camera: dict[int, list[ReplayFrame]]
    color_frame_paths: dict[int, Path]
    color_frame_images: dict[int, np.ndarray]


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description=(
            "Run the static-round FFS replay matrix, render masked RGB/PCD frame-10 boards, "
            "and export a ranked PPTX."
        )
    )
    parser.add_argument("--aligned_root", type=Path, default=ROOT / "data")
    parser.add_argument("--ffs_repo", type=Path, default=DEFAULT_FFS_REPO)
    parser.add_argument("--output_root", type=Path, default=None)
    parser.add_argument("--artifact_root", type=Path, default=None)
    parser.add_argument("--frame_idx", type=int, default=FRAME_IDX_VISUAL_DEFAULT)
    parser.add_argument("--mask_prompt", type=str, default=MASK_PROMPT_DEFAULT)
    parser.add_argument("--workspace_gib", type=int, default=8)
    parser.add_argument("--max_disp", type=int, default=DEFAULT_MAX_DISP)
    parser.add_argument("--render_width", type=int, default=848)
    parser.add_argument("--render_height", type=int, default=480)
    parser.add_argument("--reuse_artifacts", action="store_true")
    return parser.parse_args()


def format_scale_token(scale: float) -> str:
    return f"{float(scale):.1f}".replace(".", "p")


def resolve_trt_size_for_scale(scale: float) -> tuple[int, int]:
    if math.isclose(float(scale), 1.0):
        return 480, 864
    if math.isclose(float(scale), 0.75):
        return 384, 640
    if math.isclose(float(scale), 0.5):
        return 256, 448
    raise ValueError(f"Unsupported TRT scale policy: {scale}")


def build_experiment_id(*, engine: str, model_name: str, scale: float, valid_iters: int) -> str:
    return f"{engine}__model_{model_name}__scale_{format_scale_token(scale)}__iters_{int(valid_iters)}"


def build_experiment_matrix(*, ffs_repo: Path, artifacts_root: Path, max_disp: int) -> list[ExperimentConfig]:
    configs: list[ExperimentConfig] = []
    for engine in ENGINE_NAMES:
        for model_name in MODEL_NAMES:
            model_path = ffs_repo / "weights" / model_name / "model_best_bp2_serialize.pth"
            for scale in SCALE_VALUES:
                height, width = resolve_trt_size_for_scale(scale)
                for valid_iters in VALID_ITERS_VALUES:
                    experiment_id = build_experiment_id(
                        engine=engine,
                        model_name=model_name,
                        scale=scale,
                        valid_iters=valid_iters,
                    )
                    artifact_dir = (
                        artifacts_root
                        / engine
                        / f"model_{model_name}"
                        / f"scale_{format_scale_token(scale)}_iters_{int(valid_iters)}"
                    )
                    configs.append(
                        ExperimentConfig(
                            experiment_id=experiment_id,
                            engine=engine,
                            model_name=model_name,
                            model_path=str(model_path),
                            scale=float(scale),
                            valid_iters=int(valid_iters),
                            max_disp=int(max_disp),
                            engine_height=int(height),
                            engine_width=int(width),
                            artifact_dir=str(artifact_dir),
                        )
                    )
    return configs


def timestamp_token() -> str:
    return datetime.now(timezone.utc).strftime("%Y%m%d_%H%M%S")


def round_label_from_case_ref(case_ref: str) -> str:
    stem = Path(case_ref).name
    for token in ("round1", "round2", "round3"):
        if token in stem:
            return token.replace("round", "Round ")
    return stem


def load_round_case_bundle(*, aligned_root: Path, case_ref: str, frame_idx: int) -> RoundCaseBundle:
    case_dir = resolve_case_dir(aligned_root=aligned_root, case_ref=case_ref)
    metadata = load_case_metadata(case_dir)
    serial_numbers = [str(item) for item in metadata["serial_numbers"]]
    if len(serial_numbers) != 3:
        raise ValueError(f"Expected exactly 3 cameras in {case_dir}, got {len(serial_numbers)}.")

    frame_count = get_frame_count(metadata)
    if int(frame_idx) < 0 or int(frame_idx) >= frame_count:
        raise ValueError(f"frame_idx={frame_idx} out of range for {case_dir} frame_count={frame_count}.")

    calibration_reference_serials = metadata.get("calibration_reference_serials", metadata["serial_numbers"])
    c2w_list = load_calibration_transforms(
        case_dir / "calibrate.pkl",
        serial_numbers=metadata["serial_numbers"],
        calibration_reference_serials=calibration_reference_serials,
    )
    stereo_pairs_by_camera: dict[int, list[ReplayFrame]] = {}
    color_frame_paths: dict[int, Path] = {}
    color_frame_images: dict[int, np.ndarray] = {}
    for camera_idx in CAMERA_IDS:
        replay_frames: list[ReplayFrame] = []
        for frame_number in range(frame_count):
            left_path = case_dir / "ir_left" / str(camera_idx) / f"{frame_number}.png"
            right_path = case_dir / "ir_right" / str(camera_idx) / f"{frame_number}.png"
            left_image = cv2.imread(str(left_path), cv2.IMREAD_UNCHANGED)
            right_image = cv2.imread(str(right_path), cv2.IMREAD_UNCHANGED)
            if left_image is None or right_image is None:
                raise FileNotFoundError(
                    f"Missing aligned stereo pair for case={case_dir.name}, cam={camera_idx}, frame={frame_number}."
                )
            replay_frames.append(
                ReplayFrame(
                    frame_idx=int(frame_number),
                    left_image=np.asarray(left_image),
                    right_image=np.asarray(right_image),
                )
            )
        color_path = case_dir / "color" / str(camera_idx) / f"{frame_idx}.png"
        color_image = cv2.imread(str(color_path), cv2.IMREAD_COLOR)
        if color_image is None:
            raise FileNotFoundError(f"Missing RGB frame for case={case_dir.name}, cam={camera_idx}, frame={frame_idx}.")
        stereo_pairs_by_camera[int(camera_idx)] = replay_frames
        color_frame_paths[int(camera_idx)] = color_path
        color_frame_images[int(camera_idx)] = np.asarray(color_image)

    return RoundCaseBundle(
        round_label=round_label_from_case_ref(case_ref),
        case_ref=case_ref,
        case_dir=case_dir,
        metadata=metadata,
        serial_numbers=serial_numbers,
        k_ir_left_list=[np.asarray(item, dtype=np.float32) for item in metadata["K_ir_left"]],
        k_color_list=[np.asarray(item, dtype=np.float32) for item in metadata["K_color"]],
        t_ir_left_to_color_list=[np.asarray(item, dtype=np.float32) for item in metadata["T_ir_left_to_color"]],
        baselines_m=[float(item) for item in metadata["ir_baseline_m"]],
        c2w_list=[np.asarray(item, dtype=np.float32) for item in c2w_list],
        stereo_pairs_by_camera=stereo_pairs_by_camera,
        color_frame_paths=color_frame_paths,
        color_frame_images=color_frame_images,
    )


def _mask_loader_camera_clouds(bundle: RoundCaseBundle) -> list[dict[str, Any]]:
    return [
        {
            "camera_idx": int(camera_idx),
            "serial": str(bundle.serial_numbers[camera_idx]),
            "color_path": str(bundle.color_frame_paths[camera_idx]),
        }
        for camera_idx in CAMERA_IDS
    ]


def mask_cache_has_frame(*, mask_root: Path, bundle: RoundCaseBundle, frame_idx: int, text_prompt: str) -> bool:
    try:
        masks, _ = load_union_masks_for_camera_clouds(
            mask_root=mask_root,
            camera_clouds=_mask_loader_camera_clouds(bundle),
            frame_token=str(int(frame_idx)),
            text_prompt=text_prompt,
        )
    except Exception:
        return False
    return all(int(np.count_nonzero(masks.get(int(camera_idx), np.zeros((1, 1), dtype=bool)))) > 0 for camera_idx in CAMERA_IDS)


def _copy_static_mask_fallback(
    *,
    bundle: RoundCaseBundle,
    mask_root: Path,
    text_prompt: str,
    frame_idx: int,
) -> bool:
    if " ".join(parse_text_prompts(text_prompt)) != "stuffed animal":
        return False
    source_root = STATIC_MASK_FALLBACK_ROOTS.get(bundle.round_label)
    if source_root is None or not source_root.is_dir():
        return False

    if mask_root.exists():
        shutil.rmtree(mask_root, ignore_errors=True)
    shutil.copytree(source_root, mask_root)
    frame_token = str(int(frame_idx))
    for camera_idx in CAMERA_IDS:
        mask_info_path = mask_root / "mask" / f"mask_info_{int(camera_idx)}.json"
        if not mask_info_path.is_file():
            continue
        info = json.loads(mask_info_path.read_text(encoding="utf-8"))
        for obj_id in info.keys():
            obj_dir = mask_root / "mask" / str(int(camera_idx)) / str(int(obj_id))
            source_mask = obj_dir / "0.png"
            target_mask = obj_dir / f"{frame_token}.png"
            if source_mask.is_file() and not target_mask.exists():
                shutil.copy2(source_mask, target_mask)
    write_json(
        mask_root / "summary.json",
        {
            "case_root": str(bundle.case_dir),
            "output_dir": str(mask_root),
            "camera_ids": [int(item) for item in CAMERA_IDS],
            "text_prompt": str(text_prompt),
            "parsed_prompts": parse_text_prompts(text_prompt),
            "ann_frame_index": int(frame_idx),
            "generation_mode": "copied_static_frame0_fallback",
            "fallback_source_root": str(source_root),
        },
    )
    return mask_cache_has_frame(mask_root=mask_root, bundle=bundle, frame_idx=frame_idx, text_prompt=text_prompt)


def ensure_mask_cache(
    *,
    bundle: RoundCaseBundle,
    mask_root: Path,
    text_prompt: str,
    frame_idx: int,
) -> Path:
    if mask_cache_has_frame(mask_root=mask_root, bundle=bundle, frame_idx=frame_idx, text_prompt=text_prompt):
        return mask_root

    from scripts.harness.sam31_mask_helper import run_case_segmentation

    overwrite = mask_root.exists()
    try:
        run_case_segmentation(
            case_root=bundle.case_dir,
            text_prompt=text_prompt,
            camera_ids=list(CAMERA_IDS),
            output_dir=mask_root,
            source_mode="frames",
            ann_frame_index=int(frame_idx),
            overwrite=overwrite,
        )
    except Exception:
        if not _copy_static_mask_fallback(
            bundle=bundle,
            mask_root=mask_root,
            text_prompt=text_prompt,
            frame_idx=frame_idx,
        ):
            raise
    if not mask_cache_has_frame(mask_root=mask_root, bundle=bundle, frame_idx=frame_idx, text_prompt=text_prompt):
        raise RuntimeError(
            f"Mask generation did not produce non-empty frame-{frame_idx} masks for {bundle.case_dir.name}."
        )
    return mask_root


def compose_plain_matrix_board(
    *,
    image_rows: list[list[np.ndarray]],
    row_headers: list[str],
    column_headers: list[str],
    title: str | None = None,
    background_bgr: tuple[int, int, int] = (255, 255, 255),
) -> np.ndarray:
    if not image_rows or not image_rows[0]:
        raise ValueError("compose_plain_matrix_board requires a non-empty image matrix.")
    if len(image_rows) != len(row_headers):
        raise ValueError("row_headers must match image row count.")
    col_count = len(image_rows[0])
    if any(len(row) != col_count for row in image_rows):
        raise ValueError("All image rows must have the same number of columns.")
    if len(column_headers) != col_count:
        raise ValueError("column_headers must match image column count.")

    tile_h, tile_w = image_rows[0][0].shape[:2]
    for row in image_rows:
        for image in row:
            if image.shape[:2] != (tile_h, tile_w):
                raise ValueError("All board images must share the same shape.")

    padding = 16
    gap = 12
    row_label_w = 90
    col_label_h = 40
    title_h = 38 if title else 0
    board_w = row_label_w + padding * 2 + col_count * tile_w + (col_count - 1) * gap
    board_h = title_h + col_label_h + padding * 2 + len(image_rows) * tile_h + (len(image_rows) - 1) * gap
    canvas = np.full((board_h, board_w, 3), background_bgr, dtype=np.uint8)

    if title:
        cv2.putText(
            canvas,
            title,
            (padding, 26),
            cv2.FONT_HERSHEY_SIMPLEX,
            0.72,
            (0, 0, 0),
            2,
            cv2.LINE_AA,
        )

    header_y = title_h + 28
    for col_idx, header in enumerate(column_headers):
        x0 = row_label_w + padding + col_idx * (tile_w + gap)
        cv2.putText(
            canvas,
            header,
            (x0 + 8, header_y),
            cv2.FONT_HERSHEY_SIMPLEX,
            0.62,
            (0, 0, 0),
            2,
            cv2.LINE_AA,
        )

    for row_idx, row in enumerate(image_rows):
        y0 = title_h + col_label_h + padding + row_idx * (tile_h + gap)
        cv2.putText(
            canvas,
            row_headers[row_idx],
            (padding, y0 + 30),
            cv2.FONT_HERSHEY_SIMPLEX,
            0.62,
            (0, 0, 0),
            2,
            cv2.LINE_AA,
        )
        for col_idx, image in enumerate(row):
            x0 = row_label_w + padding + col_idx * (tile_w + gap)
            canvas[y0:y0 + tile_h, x0:x0 + tile_w] = image
            cv2.rectangle(canvas, (x0 - 1, y0 - 1), (x0 + tile_w, y0 + tile_h), (200, 200, 200), 1, cv2.LINE_AA)
    return canvas


def render_shared_masked_rgb_board(
    *,
    round_cases: list[RoundCaseBundle],
    mask_roots: dict[str, Path],
    text_prompt: str,
    frame_idx: int,
    output_path: Path,
) -> Path:
    image_rows: list[list[np.ndarray]] = []
    for bundle in round_cases:
        masks, _ = load_union_masks_for_camera_clouds(
            mask_root=mask_roots[bundle.round_label],
            camera_clouds=_mask_loader_camera_clouds(bundle),
            frame_token=str(int(frame_idx)),
            text_prompt=text_prompt,
        )
        row_images = [
            _mask_rgb_image(bundle.color_frame_paths[camera_idx], mask=np.asarray(masks[int(camera_idx)], dtype=bool))
            for camera_idx in CAMERA_IDS
        ]
        image_rows.append(row_images)
    board = compose_plain_matrix_board(
        image_rows=image_rows,
        row_headers=[bundle.round_label for bundle in round_cases],
        column_headers=[f"Cam {int(camera_idx) + 1}" for camera_idx in CAMERA_IDS],
        title=f"Masked RGB | frame {int(frame_idx)} | prompt={text_prompt}",
    )
    write_image(output_path, board)
    return output_path


def _fuse_camera_clouds(camera_clouds: list[dict[str, Any]]) -> tuple[np.ndarray, np.ndarray]:
    point_sets = [np.asarray(item["points"], dtype=np.float32) for item in camera_clouds if len(item["points"]) > 0]
    color_sets = [np.asarray(item["colors"], dtype=np.uint8) for item in camera_clouds if len(item["points"]) > 0]
    if not point_sets:
        return np.empty((0, 3), dtype=np.float32), np.empty((0, 3), dtype=np.uint8)
    if len(point_sets) == 1:
        return point_sets[0], color_sets[0]
    return np.concatenate(point_sets, axis=0), np.concatenate(color_sets, axis=0)


def _masked_focus_bounds(
    *,
    masked_points: np.ndarray,
    unmasked_points: np.ndarray,
) -> tuple[np.ndarray, np.ndarray]:
    if len(masked_points) >= MIN_MASKED_POINT_COUNT_FOR_FOCUS:
        bounds_min = np.asarray(masked_points, dtype=np.float32).min(axis=0)
        bounds_max = np.asarray(masked_points, dtype=np.float32).max(axis=0)
    elif len(unmasked_points) > 0:
        bounds_min = np.asarray(unmasked_points, dtype=np.float32).min(axis=0)
        bounds_max = np.asarray(unmasked_points, dtype=np.float32).max(axis=0)
    else:
        raise RuntimeError("Cannot compute focus bounds for an empty round cloud.")
    return _expand_bounds(bounds_min, bounds_max)


def _render_open3d_pinhole_offscreen(
    points: np.ndarray,
    colors: np.ndarray,
    *,
    width: int,
    height: int,
    intrinsic_matrix: np.ndarray,
    extrinsic_matrix: np.ndarray,
    point_size: float,
) -> np.ndarray:
    if len(points) == 0:
        return np.zeros((int(height), int(width), 3), dtype=np.uint8)

    import open3d as o3d
    from open3d.visualization import rendering

    renderer = rendering.OffscreenRenderer(int(width), int(height))
    material = rendering.MaterialRecord()
    material.shader = "defaultUnlit"
    material.point_size = float(point_size)
    pcd = o3d.geometry.PointCloud()
    pcd.points = o3d.utility.Vector3dVector(np.asarray(points, dtype=np.float64))
    pcd.colors = o3d.utility.Vector3dVector(np.asarray(colors[:, ::-1], dtype=np.float64) / 255.0)
    try:
        renderer.scene.set_background([0.0, 0.0, 0.0, 1.0])
        renderer.scene.add_geometry("pcd", pcd, material)
        renderer.setup_camera(
            np.asarray(intrinsic_matrix, dtype=np.float64).reshape(3, 3),
            np.asarray(extrinsic_matrix, dtype=np.float64).reshape(4, 4),
            int(width),
            int(height),
        )
        image = np.asarray(renderer.render_to_image())
    finally:
        renderer.scene.clear_geometry()
        if hasattr(renderer, "release"):
            renderer.release()
        del renderer
    return cv2.cvtColor(image, cv2.COLOR_RGBA2BGR)


def _build_round_camera_clouds(
    *,
    bundle: RoundCaseBundle,
    depth_color_by_camera: dict[int, np.ndarray],
) -> list[dict[str, Any]]:
    clouds: list[dict[str, Any]] = []
    for camera_idx in CAMERA_IDS:
        depth_color_m = np.asarray(depth_color_by_camera[int(camera_idx)], dtype=np.float32)
        color_image = bundle.color_frame_images[int(camera_idx)]
        camera_points, camera_colors, source_pixel_uv, _ = depth_to_camera_points(
            depth_color_m,
            bundle.k_color_list[int(camera_idx)],
            depth_min_m=0.2,
            depth_max_m=1.5,
            color_image=color_image,
            pixel_roi=None,
            max_points_per_camera=None,
        )
        world_points = transform_points(camera_points, bundle.c2w_list[int(camera_idx)])
        clouds.append(
            {
                "camera_idx": int(camera_idx),
                "serial": str(bundle.serial_numbers[int(camera_idx)]),
                "color_path": str(bundle.color_frame_paths[int(camera_idx)]),
                "K_color": np.asarray(bundle.k_color_list[int(camera_idx)], dtype=np.float32),
                "c2w": np.asarray(bundle.c2w_list[int(camera_idx)], dtype=np.float32),
                "points": world_points,
                "colors": camera_colors,
                "source_pixel_uv": source_pixel_uv,
                "source_camera_idx": np.full((len(world_points),), int(camera_idx), dtype=np.int16),
                "source_serial": np.full((len(world_points),), bundle.serial_numbers[int(camera_idx)], dtype=object),
            }
        )
    return clouds


def render_masked_ffs_only_pcd_board(
    *,
    round_cases: list[RoundCaseBundle],
    mask_roots: dict[str, Path],
    text_prompt: str,
    frame_idx: int,
    depth_color_by_round_camera: dict[str, dict[int, np.ndarray]],
    output_path: Path,
    render_width: int,
    render_height: int,
) -> Path:
    image_rows: list[list[np.ndarray]] = []
    for bundle in round_cases:
        camera_clouds = _build_round_camera_clouds(
            bundle=bundle,
            depth_color_by_camera=depth_color_by_round_camera[bundle.round_label],
        )
        raw_masks, _ = load_union_masks_for_camera_clouds(
            mask_root=mask_roots[bundle.round_label],
            camera_clouds=camera_clouds,
            frame_token=str(int(frame_idx)),
            text_prompt=text_prompt,
        )
        refined_masks, _ = refine_pixel_masks_with_phystwin_data_process_mask(
            camera_clouds,
            pixel_mask_by_camera=raw_masks,
        )
        masked_clouds, _ = filter_camera_clouds_with_pixel_masks(
            camera_clouds,
            pixel_mask_by_camera=refined_masks,
        )
        masked_points, masked_colors = _fuse_camera_clouds(masked_clouds)
        unmasked_points, unmasked_colors = _fuse_camera_clouds(camera_clouds)
        crop_bounds = _masked_focus_bounds(masked_points=masked_points, unmasked_points=unmasked_points)
        serial_numbers = [str(bundle.serial_numbers[int(camera_idx)]) for camera_idx in CAMERA_IDS]
        view_configs = build_original_camera_view_configs(
            c2w_list=[np.asarray(bundle.c2w_list[int(camera_idx)], dtype=np.float32) for camera_idx in CAMERA_IDS],
            serial_numbers=serial_numbers,
            look_distance=float(DEFAULT_LOOK_DISTANCE),
            camera_ids=list(range(len(CAMERA_IDS))),
        )
        row_images: list[np.ndarray] = []
        for local_camera_idx, view_config in enumerate(view_configs):
            actual_camera_idx = int(CAMERA_IDS[int(local_camera_idx)])
            source_image_size = _image_size_from_color_path(bundle.color_frame_paths[actual_camera_idx])
            target_image_size = (int(render_width), int(render_height))
            intrinsic_matrix = _scale_intrinsic_matrix(
                np.asarray(bundle.k_color_list[actual_camera_idx], dtype=np.float32),
                source_size=source_image_size,
                target_size=target_image_size,
            )
            extrinsic_matrix = np.linalg.inv(
                np.asarray(bundle.c2w_list[actual_camera_idx], dtype=np.float32).reshape(4, 4)
            ).astype(np.float32)
            cropped_points, cropped_colors = crop_points_to_bounds(masked_points, masked_colors, crop_bounds)
            render_points = cropped_points if len(cropped_points) > 0 else unmasked_points
            render_colors = cropped_colors if len(cropped_points) > 0 else unmasked_colors
            rendered = _render_open3d_pinhole_offscreen(
                render_points,
                render_colors,
                width=int(render_width),
                height=int(render_height),
                intrinsic_matrix=intrinsic_matrix,
                extrinsic_matrix=extrinsic_matrix,
                point_size=float(DEFAULT_POINT_SIZE),
            )
            row_images.append(rendered)
        image_rows.append(row_images)

    board = compose_plain_matrix_board(
        image_rows=image_rows,
        row_headers=[bundle.round_label for bundle in round_cases],
        column_headers=[f"Cam {int(camera_idx) + 1}" for camera_idx in CAMERA_IDS],
        title=(
            f"Masked FFS Point Cloud | frame {int(frame_idx)} | prompt={text_prompt} | "
            f"PhysTwin-like radius={PHYSTWIN_DATA_PROCESS_MASK_CONTRACT['radius_m']:.2f}"
        ),
    )
    write_image(output_path, board)
    return output_path


def _cleanup_runner(runner: Any) -> None:
    torch_module = getattr(runner, "torch", None)
    del runner
    gc.collect()
    if torch_module is not None and hasattr(torch_module, "cuda") and torch_module.cuda.is_available():
        torch_module.cuda.empty_cache()


def _torch_synchronize(runner: Any) -> None:
    torch_module = getattr(runner, "torch", None)
    if torch_module is not None and hasattr(torch_module, "cuda") and torch_module.cuda.is_available():
        torch_module.cuda.synchronize()


def _run_logged_subprocess(*, cmd: list[str], log_path: Path) -> tuple[int, str]:
    log_path.parent.mkdir(parents=True, exist_ok=True)
    completed = subprocess.run(
        cmd,
        cwd=ROOT,
        capture_output=True,
        text=True,
    )
    combined = "\n".join(
        [
            f"$ {' '.join(cmd)}",
            "",
            completed.stdout.strip(),
            completed.stderr.strip(),
        ]
    ).strip() + "\n"
    log_path.write_text(combined, encoding="utf-8")
    return int(completed.returncode), combined


def _two_stage_artifact_ready(artifact_dir: Path) -> bool:
    return (
        (artifact_dir / "feature_runner.engine").is_file()
        and (artifact_dir / "post_runner.engine").is_file()
        and (artifact_dir / "onnx.yaml").is_file()
    )


def _single_engine_artifact_ready(artifact_dir: Path) -> bool:
    return (
        any(path.is_file() for path in artifact_dir.glob("*.engine"))
        and any(path.is_file() for path in artifact_dir.glob("*.yaml"))
    )


def _validate_existing_artifact(config: ExperimentConfig) -> bool:
    artifact_dir = Path(config.artifact_dir)
    if config.engine == "two_stage_fp16":
        if not _two_stage_artifact_ready(artifact_dir):
            return False
        batch_size = resolve_tensorrt_engine_static_batch_size(trt_mode="two_stage", model_dir=artifact_dir)
        return int(batch_size) == 1
    if config.engine == "single_engine_fp32":
        if not _single_engine_artifact_ready(artifact_dir):
            return False
        batch_size = resolve_tensorrt_engine_static_batch_size(trt_mode="single_engine", model_dir=artifact_dir)
        return int(batch_size) == 1
    raise ValueError(f"Unsupported engine: {config.engine}")


def ensure_trt_artifact(
    *,
    config: ExperimentConfig,
    ffs_repo: Path,
    workspace_gib: int,
    reuse_artifacts: bool,
    experiment_dir: Path,
) -> tuple[bool, Path, str | None]:
    artifact_dir = Path(config.artifact_dir)
    artifact_dir.mkdir(parents=True, exist_ok=True)
    build_log = experiment_dir / "artifact_build.log"
    if reuse_artifacts and _validate_existing_artifact(config):
        build_log.write_text("Reused existing artifact.\n", encoding="utf-8")
        return True, build_log, None

    script_path: Path
    cmd: list[str]
    common = [
        sys.executable,
        str(ROOT / "scripts" / "harness"),
    ]
    if config.engine == "two_stage_fp16":
        script_path = ROOT / "scripts" / "harness" / "verify_ffs_tensorrt_wsl.py"
        cmd = [
            sys.executable,
            str(script_path),
            "--ffs_repo",
            str(ffs_repo),
            "--model_path",
            str(config.model_path),
            "--out_dir",
            str(artifact_dir),
            "--height",
            str(config.engine_height),
            "--width",
            str(config.engine_width),
            "--batch_size",
            "1",
            "--valid_iters",
            str(config.valid_iters),
            "--max_disp",
            str(config.max_disp),
            "--workspace_gib",
            str(workspace_gib),
            "--skip_profiles",
        ]
    elif config.engine == "single_engine_fp32":
        script_path = ROOT / "scripts" / "harness" / "verify_ffs_single_engine_tensorrt_wsl.py"
        cmd = [
            sys.executable,
            str(script_path),
            "--ffs_repo",
            str(ffs_repo),
            "--model_path",
            str(config.model_path),
            "--out_dir",
            str(artifact_dir),
            "--height",
            str(config.engine_height),
            "--width",
            str(config.engine_width),
            "--batch_size",
            "1",
            "--valid_iters",
            str(config.valid_iters),
            "--max_disp",
            str(config.max_disp),
            "--workspace_gib",
            str(workspace_gib),
        ]
    else:
        raise ValueError(f"Unsupported engine: {config.engine}")

    returncode, combined_output = _run_logged_subprocess(cmd=cmd, log_path=build_log)
    if returncode != 0:
        return False, build_log, combined_output.strip() or f"Artifact build failed with exit={returncode}."
    if not _validate_existing_artifact(config):
        return False, build_log, "Artifact build completed but validation did not find a batch=1 runnable engine."
    return True, build_log, None


def build_runner_for_experiment(config: ExperimentConfig, *, ffs_repo: Path) -> Any:
    artifact_dir = Path(config.artifact_dir)
    if config.engine == "two_stage_fp16":
        return FastFoundationStereoTensorRTRunner(ffs_repo=ffs_repo, model_dir=artifact_dir)
    if config.engine == "single_engine_fp32":
        return FastFoundationStereoSingleEngineTensorRTRunner(ffs_repo=ffs_repo, model_dir=artifact_dir)
    raise ValueError(f"Unsupported engine: {config.engine}")


def load_replay_frames_for_camera(*, case_dir: Path, camera_idx: int) -> list[ReplayFrame]:
    left_dir = case_dir / "ir_left" / str(int(camera_idx))
    right_dir = case_dir / "ir_right" / str(int(camera_idx))
    if not left_dir.is_dir() or not right_dir.is_dir():
        raise FileNotFoundError(f"Missing aligned stereo directories for case={case_dir.name}, cam={camera_idx}.")
    frame_ids = sorted(int(path.stem) for path in left_dir.glob("*.png"))
    replay_frames: list[ReplayFrame] = []
    for frame_number in frame_ids:
        left_path = left_dir / f"{frame_number}.png"
        right_path = right_dir / f"{frame_number}.png"
        left_image = cv2.imread(str(left_path), cv2.IMREAD_UNCHANGED)
        right_image = cv2.imread(str(right_path), cv2.IMREAD_UNCHANGED)
        if left_image is None or right_image is None:
            raise FileNotFoundError(
                f"Missing aligned stereo pair for case={case_dir.name}, cam={camera_idx}, frame={frame_number}."
            )
        replay_frames.append(
            ReplayFrame(
                frame_idx=int(frame_number),
                left_image=np.asarray(left_image),
                right_image=np.asarray(right_image),
            )
        )
    if not replay_frames:
        raise ValueError(f"No replay frames found for case={case_dir.name}, cam={camera_idx}.")
    return replay_frames


def _benchmark_round_camera_job(job: RoundCameraBenchmarkJob) -> RoundCameraBenchmarkResult:
    runner = build_runner_for_experiment(job.config, ffs_repo=Path(job.ffs_repo))
    try:
        replay_frames = load_replay_frames_for_camera(case_dir=Path(job.case_dir), camera_idx=int(job.camera_idx))
        depth_color_frame: np.ndarray | None = None
        for warmup_frame in replay_frames[:10]:
            runner.run_pair(
                warmup_frame.left_image,
                warmup_frame.right_image,
                K_ir_left=np.asarray(job.k_ir_left, dtype=np.float32),
                baseline_m=float(job.baseline_m),
                audit_mode=False,
            )
        _torch_synchronize(runner)
        start_s = time.perf_counter()
        for replay_frame in replay_frames:
            output = runner.run_pair(
                replay_frame.left_image,
                replay_frame.right_image,
                K_ir_left=np.asarray(job.k_ir_left, dtype=np.float32),
                baseline_m=float(job.baseline_m),
                audit_mode=False,
            )
            if int(replay_frame.frame_idx) == int(job.frame_idx):
                depth_color_frame = align_depth_to_color(
                    np.asarray(output["depth_ir_left_m"], dtype=np.float32),
                    np.asarray(output["K_ir_left_used"], dtype=np.float32),
                    np.asarray(job.t_ir_left_to_color, dtype=np.float32),
                    np.asarray(job.k_color, dtype=np.float32),
                    output_shape=tuple(int(item) for item in job.color_output_shape),
                )
        _torch_synchronize(runner)
        elapsed_s = max(1e-9, float(time.perf_counter() - start_s))
        if depth_color_frame is None:
            raise RuntimeError(
                f"Did not capture frame_idx={int(job.frame_idx)} while benchmarking {job.round_label} cam={int(job.camera_idx)}."
            )
        return RoundCameraBenchmarkResult(
            round_label=str(job.round_label),
            camera_idx=int(job.camera_idx),
            fps=float(len(replay_frames) / elapsed_s),
            depth_color_m=np.asarray(depth_color_frame, dtype=np.float32),
        )
    finally:
        _cleanup_runner(runner)


def _benchmark_round_camera_job_entry(*, job: RoundCameraBenchmarkJob, result_queue: Any) -> None:
    try:
        result = _benchmark_round_camera_job(job)
        result_queue.put(
            {
                "round_label": str(result.round_label),
                "camera_idx": int(result.camera_idx),
                "fps": float(result.fps),
                "depth_color_m": np.asarray(result.depth_color_m, dtype=np.float32),
            }
        )
    except Exception as exc:
        result_queue.put(
            {
                "round_label": str(job.round_label),
                "camera_idx": int(job.camera_idx),
                "error": f"{type(exc).__name__}: {exc}",
            }
        )


def benchmark_experiment(
    *,
    config: ExperimentConfig,
    round_cases: list[RoundCaseBundle],
    ffs_repo: Path,
    frame_idx: int,
) -> tuple[dict[str, dict[int, float]], dict[str, dict[int, np.ndarray]]]:
    fps_by_round: dict[str, dict[int, float]] = {}
    frame_depths: dict[str, dict[int, np.ndarray]] = {}
    mp_ctx = mp.get_context("spawn")
    for bundle in round_cases:
        round_jobs = [
            RoundCameraBenchmarkJob(
                config=config,
                ffs_repo=str(ffs_repo),
                round_label=str(bundle.round_label),
                case_dir=str(bundle.case_dir),
                camera_idx=int(camera_idx),
                frame_idx=int(frame_idx),
                k_ir_left=np.asarray(bundle.k_ir_left_list[int(camera_idx)], dtype=np.float32),
                k_color=np.asarray(bundle.k_color_list[int(camera_idx)], dtype=np.float32),
                t_ir_left_to_color=np.asarray(bundle.t_ir_left_to_color_list[int(camera_idx)], dtype=np.float32),
                baseline_m=float(bundle.baselines_m[int(camera_idx)]),
                color_output_shape=(
                    int(bundle.color_frame_images[int(camera_idx)].shape[0]),
                    int(bundle.color_frame_images[int(camera_idx)].shape[1]),
                ),
            )
            for camera_idx in CAMERA_IDS
        ]
        result_queue = mp_ctx.Queue()
        workers = [
            mp_ctx.Process(
                target=_benchmark_round_camera_job_entry,
                kwargs={"job": job, "result_queue": result_queue},
            )
            for job in round_jobs
        ]
        for worker in workers:
            worker.start()

        round_fps: dict[int, float] = {}
        round_depths: dict[int, np.ndarray] = {}
        errors: list[str] = []
        try:
            for _ in round_jobs:
                payload = result_queue.get(timeout=float(ROUND_BENCHMARK_RESULT_TIMEOUT_S))
                camera_idx = int(payload["camera_idx"])
                if "error" in payload:
                    errors.append(f"{bundle.round_label} cam{int(camera_idx) + 1}: {payload['error']}")
                    continue
                round_fps[camera_idx] = float(payload["fps"])
                round_depths[camera_idx] = np.asarray(payload["depth_color_m"], dtype=np.float32)
        except queue.Empty as exc:
            errors.append(
                f"{bundle.round_label}: timed out waiting for concurrent round benchmark worker results after "
                f"{float(ROUND_BENCHMARK_RESULT_TIMEOUT_S):.1f}s."
            )
        finally:
            for worker in workers:
                worker.join(timeout=2.0)
            for worker in workers:
                if worker.is_alive():
                    worker.terminate()
                    worker.join(timeout=2.0)
            try:
                result_queue.close()
            except Exception:
                pass
        if errors:
            raise RuntimeError("; ".join(errors))
        fps_by_round[bundle.round_label] = round_fps
        frame_depths[bundle.round_label] = round_depths
    return fps_by_round, frame_depths


def flatten_fps_values(fps_by_round: dict[str, dict[int, float]]) -> list[float]:
    values: list[float] = []
    for round_label in sorted(fps_by_round):
        for camera_idx in CAMERA_IDS:
            values.append(float(fps_by_round[round_label][int(camera_idx)]))
    return values


def overall_mean_fps(fps_by_round: dict[str, dict[int, float]]) -> float:
    values = flatten_fps_values(fps_by_round)
    return float(sum(values) / max(1, len(values)))


def build_results_row(
    *,
    config: ExperimentConfig,
    fps_by_round: dict[str, dict[int, float]],
    rgb_board_path: Path,
    pcd_board_path: Path,
    experiment_dir: Path,
) -> dict[str, Any]:
    row: dict[str, Any] = {
        "experiment_id": config.experiment_id,
        "engine": config.engine,
        "model_name": config.model_name,
        "scale": float(config.scale),
        "valid_iters": int(config.valid_iters),
        "engine_height": int(config.engine_height),
        "engine_width": int(config.engine_width),
        "artifact_dir": str(config.artifact_dir),
        "rgb_board_path": str(rgb_board_path),
        "pcd_board_path": str(pcd_board_path),
        "experiment_dir": str(experiment_dir),
    }
    for round_label in sorted(fps_by_round):
        round_token = round_label.lower().replace(" ", "")
        for camera_idx in CAMERA_IDS:
            row[f"{round_token}_cam{int(camera_idx) + 1}_fps"] = float(fps_by_round[round_label][int(camera_idx)])
    row["overall_mean_fps"] = overall_mean_fps(fps_by_round)
    return row


def build_results_row_from_summary(summary: dict[str, Any]) -> dict[str, Any]:
    config = dict(summary["config"])
    fps_by_round = {
        str(round_label): {int(camera_idx): float(value) for camera_idx, value in round_values.items()}
        for round_label, round_values in dict(summary["fps_by_round"]).items()
    }
    experiment_dir = Path(str(summary["pcd_board_path"])).resolve().parent
    return build_results_row(
        config=ExperimentConfig(
            experiment_id=str(config["experiment_id"]),
            engine=str(config["engine"]),
            model_name=str(config["model_name"]),
            model_path=str(config["model_path"]),
            scale=float(config["scale"]),
            valid_iters=int(config["valid_iters"]),
            max_disp=int(config["max_disp"]),
            engine_height=int(config["engine_height"]),
            engine_width=int(config["engine_width"]),
            artifact_dir=str(config["artifact_dir"]),
        ),
        fps_by_round=fps_by_round,
        rgb_board_path=Path(str(summary["rgb_board_path"])),
        pcd_board_path=Path(str(summary["pcd_board_path"])),
        experiment_dir=experiment_dir,
    )


def write_results_csv(*, csv_path: Path, rows: list[dict[str, Any]]) -> None:
    csv_path.parent.mkdir(parents=True, exist_ok=True)
    if not rows:
        csv_path.write_text("", encoding="utf-8")
        return
    fieldnames = list(rows[0].keys())
    with csv_path.open("w", encoding="utf-8", newline="") as handle:
        writer = csv.DictWriter(handle, fieldnames=fieldnames)
        writer.writeheader()
        for row in rows:
            writer.writerow(row)


def ensure_python_pptx_available() -> None:
    try:
        import pptx  # noqa: F401
    except Exception as exc:  # pragma: no cover - exercised in manual/full runs
        raise RuntimeError(
            "python-pptx is required for PPT export. Install it in the current environment before rerunning."
        ) from exc


def _add_white_background(slide: Any) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb_color(255, 255, 255)


def _rgb_color(r: int, g: int, b: int):
    from pptx.dml.color import RGBColor

    return RGBColor(int(r), int(g), int(b))


def _add_textbox(
    *,
    slide: Any,
    left: int,
    top: int,
    width: int,
    height: int,
    lines: list[str],
    font_size_pt: float,
    bold_first: bool = False,
) -> None:
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt

    text_box = slide.shapes.add_textbox(left, top, width, height)
    frame = text_box.text_frame
    frame.clear()
    frame.word_wrap = True
    for line_idx, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if line_idx == 0 else frame.add_paragraph()
        paragraph.text = str(line)
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(1)
        for run in paragraph.runs:
            run.font.size = Pt(font_size_pt)
            run.font.name = "Arial"
            run.font.bold = bool(bold_first and line_idx == 0)
            run.font.color.rgb = _rgb_color(0, 0, 0)


def _add_fitted_picture(
    *,
    slide: Any,
    image_path: Path,
    slide_width: int,
    slide_height: int,
    top_margin_in: float = 0.65,
    bottom_margin_in: float = 0.25,
    side_margin_in: float = 0.35,
) -> None:
    from pptx.util import Inches

    image = cv2.imread(str(image_path), cv2.IMREAD_COLOR)
    if image is None:
        raise FileNotFoundError(f"Missing PPT image: {image_path}")
    slide_w = int(slide_width)
    slide_h = int(slide_height)
    top_margin = int(Inches(top_margin_in))
    bottom_margin = int(Inches(bottom_margin_in))
    side_margin = int(Inches(side_margin_in))
    max_w = slide_w - side_margin * 2
    max_h = slide_h - top_margin - bottom_margin
    img_h, img_w = image.shape[:2]
    scale = min(float(max_w) / float(img_w), float(max_h) / float(img_h))
    width = int(round(img_w * scale))
    height = int(round(img_h * scale))
    left = int((slide_w - width) / 2)
    top = int(top_margin + (max_h - height) / 2)
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def export_pptx(
    *,
    pptx_path: Path,
    manifest: dict[str, Any],
    successful_rows: list[dict[str, Any]],
    failures: list[ExperimentFailure],
    shared_rgb_board_path: Path | None = None,
) -> None:
    ensure_python_pptx_available()
    from pptx import Presentation
    from pptx.util import Inches

    pptx_path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for row in successful_rows:
        slide = prs.slides.add_slide(blank_layout)
        _add_white_background(slide)
        summary_lines = [
            f"engine={row['engine']} | model={row['model_name']} | scale={row['scale']} | valid_iters={row['valid_iters']} | overall_mean_fps={row['overall_mean_fps']:.2f}",
            f"Round 1 FPS | cam1={row['round1_cam1_fps']:.2f} | cam2={row['round1_cam2_fps']:.2f} | cam3={row['round1_cam3_fps']:.2f}",
            f"Round 2 FPS | cam1={row['round2_cam1_fps']:.2f} | cam2={row['round2_cam2_fps']:.2f} | cam3={row['round2_cam3_fps']:.2f}",
            f"Round 3 FPS | cam1={row['round3_cam1_fps']:.2f} | cam2={row['round3_cam2_fps']:.2f} | cam3={row['round3_cam3_fps']:.2f}",
        ]
        _add_textbox(
            slide=slide,
            left=Inches(0.45),
            top=Inches(0.12),
            width=Inches(12.4),
            height=Inches(0.72),
            lines=summary_lines,
            font_size_pt=9.5,
            bold_first=False,
        )
        _add_fitted_picture(
            slide=slide,
            image_path=Path(row["pcd_board_path"]),
            slide_width=prs.slide_width,
            slide_height=prs.slide_height,
            top_margin_in=0.92,
            bottom_margin_in=0.22,
            side_margin_in=0.32,
        )

    prs.save(str(pptx_path))


def run_experiment_matrix(args: argparse.Namespace) -> dict[str, Any]:
    aligned_root = Path(args.aligned_root).resolve()
    ffs_repo = Path(args.ffs_repo).resolve()
    if not ffs_repo.exists():
        raise FileNotFoundError(f"FFS repo not found: {ffs_repo}")

    output_root = (
        Path(args.output_root).resolve()
        if args.output_root is not None
        else (ROOT / "data" / "experiments" / f"ffs_static_replay_matrix_{timestamp_token()}").resolve()
    )
    artifacts_root = (
        Path(args.artifact_root).resolve()
        if args.artifact_root is not None
        else (output_root / "artifacts").resolve()
    )
    experiments_root = output_root / "experiments"
    mask_cache_root = output_root / "mask_cache"
    ppt_root = output_root / "ppt"
    output_root.mkdir(parents=True, exist_ok=True)
    artifacts_root.mkdir(parents=True, exist_ok=True)
    experiments_root.mkdir(parents=True, exist_ok=True)
    mask_cache_root.mkdir(parents=True, exist_ok=True)
    ppt_root.mkdir(parents=True, exist_ok=True)

    configs = build_experiment_matrix(
        ffs_repo=ffs_repo,
        artifacts_root=artifacts_root,
        max_disp=int(args.max_disp),
    )
    round_cases = [
        load_round_case_bundle(aligned_root=aligned_root, case_ref=case_ref, frame_idx=int(args.frame_idx))
        for case_ref in ROUND_CASE_REFS
    ]

    mask_roots: dict[str, Path] = {}
    for bundle in round_cases:
        mask_root = mask_cache_root / bundle.case_dir.name
        ensure_mask_cache(
            bundle=bundle,
            mask_root=mask_root,
            text_prompt=str(args.mask_prompt),
            frame_idx=int(args.frame_idx),
        )
        mask_roots[bundle.round_label] = mask_root

    shared_rgb_board_path = output_root / "shared_masked_rgb_board.png"
    render_shared_masked_rgb_board(
        round_cases=round_cases,
        mask_roots=mask_roots,
        text_prompt=str(args.mask_prompt),
        frame_idx=int(args.frame_idx),
        output_path=shared_rgb_board_path,
    )

    success_rows: list[dict[str, Any]] = []
    failures: list[ExperimentFailure] = []
    for config in configs:
        experiment_dir = experiments_root / config.experiment_id
        experiment_dir.mkdir(parents=True, exist_ok=True)
        success, build_log, build_error = ensure_trt_artifact(
            config=config,
            ffs_repo=ffs_repo,
            workspace_gib=int(args.workspace_gib),
            reuse_artifacts=bool(args.reuse_artifacts),
            experiment_dir=experiment_dir,
        )
        if not success:
            failure = ExperimentFailure(
                experiment_id=config.experiment_id,
                engine=config.engine,
                model_name=config.model_name,
                scale=float(config.scale),
                valid_iters=int(config.valid_iters),
                stage="artifact_build",
                message=str(build_error or "artifact build failed"),
                log_path=str(build_log),
            )
            failures.append(failure)
            write_json(experiment_dir / "failure.json", asdict(failure))
            continue

        try:
            fps_by_round, frame_depths = benchmark_experiment(
                config=config,
                round_cases=round_cases,
                ffs_repo=ffs_repo,
                frame_idx=int(args.frame_idx),
            )
            pcd_board_path = experiment_dir / "pcd_board.png"
            render_masked_ffs_only_pcd_board(
                round_cases=round_cases,
                mask_roots=mask_roots,
                text_prompt=str(args.mask_prompt),
                frame_idx=int(args.frame_idx),
                depth_color_by_round_camera=frame_depths,
                output_path=pcd_board_path,
                render_width=int(args.render_width),
                render_height=int(args.render_height),
            )
            result_row = build_results_row(
                config=config,
                fps_by_round=fps_by_round,
                rgb_board_path=shared_rgb_board_path,
                pcd_board_path=pcd_board_path,
                experiment_dir=experiment_dir,
            )
            success_rows.append(result_row)
            write_json(
                experiment_dir / "summary.json",
                {
                    "config": asdict(config),
                    "fps_by_round": {
                        round_label: {str(camera_idx): float(fps) for camera_idx, fps in round_fps.items()}
                        for round_label, round_fps in fps_by_round.items()
                    },
                    "overall_mean_fps": float(result_row["overall_mean_fps"]),
                    "rgb_board_path": str(shared_rgb_board_path),
                    "pcd_board_path": str(pcd_board_path),
                    "artifact_build_log": str(build_log),
                },
            )
        except Exception as exc:  # pragma: no cover - manual/full-run path
            failure = ExperimentFailure(
                experiment_id=config.experiment_id,
                engine=config.engine,
                model_name=config.model_name,
                scale=float(config.scale),
                valid_iters=int(config.valid_iters),
                stage="benchmark_or_render",
                message=str(exc),
                log_path=None,
            )
            failures.append(failure)
            write_json(experiment_dir / "failure.json", asdict(failure))

    success_rows = sorted(success_rows, key=lambda row: float(row["overall_mean_fps"]), reverse=True)
    results_csv_path = output_root / "results.csv"
    write_results_csv(csv_path=results_csv_path, rows=success_rows)

    manifest = {
        "generated_at_utc": datetime.now(timezone.utc).isoformat(),
        "output_root": str(output_root),
        "artifact_root": str(artifacts_root),
        "aligned_root": str(aligned_root),
        "ffs_repo": str(ffs_repo),
        "frame_idx": int(args.frame_idx),
        "mask_prompt": str(args.mask_prompt),
        "parsed_prompts": parse_text_prompts(str(args.mask_prompt)),
        "configs": [asdict(config) for config in configs],
        "successful_experiment_ids": [row["experiment_id"] for row in success_rows],
        "failures": [asdict(item) for item in failures],
        "shared_rgb_board_path": str(shared_rgb_board_path),
        "results_csv_path": str(results_csv_path),
    }
    write_json(output_root / "manifest.json", manifest)

    pptx_path = ppt_root / "ffs_static_replay_matrix.pptx"
    export_pptx(
        pptx_path=pptx_path,
        manifest=manifest,
        successful_rows=success_rows,
        failures=failures,
        shared_rgb_board_path=shared_rgb_board_path,
    )
    manifest["pptx_path"] = str(pptx_path)
    write_json(output_root / "manifest.json", manifest)
    return manifest


def main() -> int:
    manifest = run_experiment_matrix(parse_args())
    print(f"Static replay matrix written to {manifest['output_root']}")
    print(f"PPTX: {manifest['pptx_path']}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
